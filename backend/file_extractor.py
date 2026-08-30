"""File text extraction for the AI Assistant chatbot.

Supports .pdf, .doc/.docx, .csv, .txt, .xlsx, .pptx and image (.jpg/.jpeg/
.png) uploads:

  * resolve_file_type()   — validate a file by extension AND magic bytes
                            (never trust the extension alone)
  * extract_text()        — pull readable text out per file type with caps
                            (20 PDF pages, 100 CSV rows, 200 XLSX rows,
                            30 PPTX slides, 3000 words)

Scanned/image-only PDFs, corrupt files and empty files return clear,
user-facing error messages (raises FileExtractionError). Images are never
OCR'd/analyzed — extract_text() raises FileExtractionError with a
user-facing explanation instead, same as any other unreadable attachment.
"""

import csv
import io
import re

MAX_PDF_PAGES = 20
MAX_CSV_ROWS = 100
MAX_XLSX_ROWS = 200
MAX_PPTX_SLIDES = 30
MAX_WORDS = 3000

SCANNED_PDF_MESSAGE = (
    "I couldn't find any readable text in this PDF. This can happen with a "
    "scanned/image-only PDF, or if a printed/exported page was saved before "
    "it finished loading (its pages come out blank). If this was a "
    "browser-printed report, please make sure the page is fully loaded "
    "before printing, then re-save and re-upload it."
)
CORRUPT_FILE_MESSAGE = (
    "Sorry, I couldn't read this file. It may be corrupted or in an "
    "unsupported format."
)
EMPTY_FILE_MESSAGE = "The file is empty."
EMPTY_PPTX_MESSAGE = "This presentation appears to have no readable text."
IMAGE_MESSAGE = (
    "I can see you've attached an image, but I can't currently read image "
    "content. Please describe what's in the image, or attach a text-based "
    "file instead."
)
TRUNCATION_NOTE = "[Content truncated - file is longer than shown]"


class FileExtractionError(Exception):
    """Raised when a file cannot be read; message is user-facing."""


# ---------------------------------------------------------------------------
# magic-byte checks
# ---------------------------------------------------------------------------

_PDF_MAGIC = b"%PDF-"
_ZIP_MAGIC = b"PK\x03\x04"
_OLE2_MAGIC = b"\xD0\xCF\x11\xE0"
_PNG_MAGIC = b"\x89PNG"
_JPEG_MAGIC = b"\xFF\xD8\xFF"

# Known binary magic sequences that must NOT appear in .csv/.txt content.
_BINARY_MAGICS = (
    b"%PDF-", b"PK\x03\x04", b"\xD0\xCF\x11\xE0", b"\x89PNG", b"\xFF\xD8\xFF",
    b"GIF8", b"BM", b"MZ", b"\x7fELF", b"\x00\x00\x00\x18ftyp", b"RIFF",
)

EXTENSION_TO_TYPE = {
    ".pdf": "pdf",
    ".doc": "doc",
    ".docx": "docx",
    ".csv": "csv",
    ".txt": "txt",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".jpg": "jpg",
    ".jpeg": "jpeg",
    ".png": "png",
}


def _head(path, n=32):
    with open(path, "rb") as handle:
        return handle.read(n)


def _looks_binary(data: bytes) -> bool:
    """A file is treated as binary when it contains NUL bytes or known magic."""
    if any(magic in data[: len(magic)] for magic in _BINARY_MAGICS):
        return True
    sample = data[:4096]
    if not sample:
        return False
    return sample.count(b"\x00") > 0


def resolve_file_type(filename: str, path: str):
    """Validate a file by extension AND magic bytes.

    Returns the canonical type string ('pdf' | 'doc' | 'docx' | 'csv' | 'txt').
    Raises FileExtractionError for an unsupported/mismatched file.
    """
    ext = (str(filename or "").rsplit(".", 1)[-1].lower() if "." in str(filename) else "")
    ext = f".{ext}"
    declared = EXTENSION_TO_TYPE.get(ext)
    if declared is None:
        raise FileExtractionError(
            "Unsupported file type. Please upload PDF, Word, Excel, PowerPoint, "
            "CSV, TXT, or image (JPG/PNG) files."
        )

    head = _head(path)
    if not head:
        raise FileExtractionError(EMPTY_FILE_MESSAGE)

    if declared == "pdf" and not head.startswith(_PDF_MAGIC):
        raise FileExtractionError(
            "Unsupported file type. The file does not appear to be a valid PDF."
        )
    if declared == "docx" and not head.startswith(_ZIP_MAGIC):
        raise FileExtractionError(
            "Unsupported file type. The file does not appear to be a valid Word (.docx) document."
        )
    if declared == "doc" and not head.startswith(_OLE2_MAGIC):
        raise FileExtractionError(
            "Unsupported file type. The file does not appear to be a valid Word (.doc) document."
        )
    if declared == "xlsx" and not head.startswith(_ZIP_MAGIC):
        raise FileExtractionError(
            "Unsupported file type. The file does not appear to be a valid Excel (.xlsx) file."
        )
    if declared == "pptx" and not head.startswith(_ZIP_MAGIC):
        raise FileExtractionError(
            "Unsupported file type. The file does not appear to be a valid PowerPoint (.pptx) file."
        )
    if declared in ("jpg", "jpeg") and not head.startswith(_JPEG_MAGIC):
        raise FileExtractionError(
            "Unsupported file type. The file does not appear to be a valid JPEG image."
        )
    if declared == "png" and not head.startswith(_PNG_MAGIC):
        raise FileExtractionError(
            "Unsupported file type. The file does not appear to be a valid PNG image."
        )
    if declared in ("csv", "txt") and _looks_binary(head):
        raise FileExtractionError(
            "Unsupported file type. Please upload PDF, Word, CSV, or TXT files."
        )
    return declared


# ---------------------------------------------------------------------------
# extraction
# ---------------------------------------------------------------------------

def _extract_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    try:
        reader = PdfReader(path)
    except Exception as exc:
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    pages = reader.pages[:MAX_PDF_PAGES]
    page_texts = []
    for page in pages:
        try:
            page_texts.append((page.extract_text() or "").strip())
        except Exception:
            page_texts.append("")

    # A page with no real text still produces a non-empty "--- page N ---"
    # marker once formatted, which used to make the old `if not text` check
    # pass even when every page was blank (observed on real uploads: browser
    # print-to-PDF output whose page content streams are empty at the PDF
    # level -- 0 bytes, confirmed by inspecting them directly -- extracted
    # as 46 chars of pure page markers and was returned as if it were real
    # content). Check the actual page text for emptiness BEFORE adding markers.
    if not any(page_texts):
        raise FileExtractionError(SCANNED_PDF_MESSAGE)

    parts = [f"--- page {i} ---\n{t}" for i, t in enumerate(page_texts, start=1)]
    return "\n".join(parts).strip()


def _extract_docx(path: str) -> str:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    try:
        document = Document(path)
    except Exception as exc:
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    parts = []
    for paragraph in document.paragraphs:
        if paragraph.text and paragraph.text.strip():
            parts.append(paragraph.text.strip())
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_doc(path: str) -> str:
    """Legacy .doc (OLE2 binary) fallback: python-docx cannot read these, so
    extract human-readable strings directly."""
    try:
        with open(path, "rb") as handle:
            data = handle.read()
    except OSError as exc:
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc
    # Readable ASCII runs of >= 4 chars (word-windowed at UTF-8/ASCII).
    chunks = re.findall(rb"[\x20-\x7e]{4,}", data)
    text = "\n".join(c.decode("latin-1") for c in chunks).strip()
    if not text:
        raise FileExtractionError(
            "Sorry, I couldn't read this .doc file. Please save it as .docx or .txt and try again."
        )
    return text


def _extract_csv(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="replace", newline="") as handle:
            rows = list(csv.reader(handle))
    except OSError as exc:
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    if not rows:
        raise FileExtractionError(EMPTY_FILE_MESSAGE)

    total = len(rows)
    shown = rows[:MAX_CSV_ROWS]
    widths = [
        max(len(str(row[i])) if i < len(row) else 0 for row in shown)
        for i in range(max((len(row) for row in shown), default=0))
    ]

    def fmt(row):
        cells = [str(c) for c in row]
        padded = [cells[i].ljust(widths[i]) if i < len(widths) else cells[i] for i in range(len(cells))]
        return " | ".join(padded).rstrip()

    lines = [fmt(row) for row in shown]
    if total > MAX_CSV_ROWS:
        lines.append(f"... (showing first {MAX_CSV_ROWS} of {total} rows)")
    return "\n".join(lines)


def _extract_xlsx(path: str) -> str:
    try:
        import pandas as pd
    except ImportError as exc:  # pragma: no cover
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    try:
        sheets = pd.read_excel(path, sheet_name=None, engine="openpyxl")
    except Exception as exc:
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    if not sheets:
        raise FileExtractionError(EMPTY_FILE_MESSAGE)

    parts = []
    rows_used = 0
    truncated = False
    for sheet_name, df in sheets.items():
        if rows_used >= MAX_XLSX_ROWS:
            truncated = True
            break
        remaining = MAX_XLSX_ROWS - rows_used
        shown = df.head(remaining)
        rows_used += len(shown)
        if len(df) > len(shown):
            truncated = True
        if shown.empty:
            continue
        parts.append(f"--- sheet: {sheet_name} ---\n{shown.to_string(index=False)}")

    text = "\n\n".join(parts).strip()
    if not text:
        raise FileExtractionError(EMPTY_FILE_MESSAGE)
    if truncated:
        text += f"\n\n[Content truncated - showing first {MAX_XLSX_ROWS} rows across all sheets]"
    return text


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    try:
        presentation = Presentation(path)
    except Exception as exc:
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc

    slides = list(presentation.slides)
    total = len(slides)
    shown = slides[:MAX_PPTX_SLIDES]

    parts = []
    for index, slide in enumerate(shown, start=1):
        lines = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for paragraph in shape.text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs).strip()
                if line:
                    lines.append(line)
        if lines:
            parts.append(f"--- slide {index} ---\n" + "\n".join(lines))

    if total > MAX_PPTX_SLIDES:
        parts.append(f"... (showing first {MAX_PPTX_SLIDES} of {total} slides)")

    text = "\n\n".join(parts).strip()
    if not text:
        raise FileExtractionError(EMPTY_PPTX_MESSAGE)
    return text


def _extract_txt(path: str) -> str:
    raw = None
    try:
        with open(path, "r", encoding="utf-8") as handle:
            raw = handle.read()
    except UnicodeDecodeError:
        try:
            with open(path, "r", encoding="latin-1") as handle:
                raw = handle.read()
        except OSError as exc:
            raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc
    except OSError as exc:
        raise FileExtractionError(CORRUPT_FILE_MESSAGE) from exc
    if raw is None:
        raise FileExtractionError(EMPTY_FILE_MESSAGE)
    return raw.strip()


def _cap_words(text: str) -> str:
    """Cap extracted text at MAX_WORDS; append a truncation note when cut."""
    if not text:
        return text
    words = text.split()
    if len(words) <= MAX_WORDS:
        return text
    trimmed = " ".join(words[:MAX_WORDS]).rstrip(" \t-|")
    return f"{trimmed}\n\n{TRUNCATION_NOTE}"


def extract_text(path: str, file_type: str) -> str:
    """Extract readable text from a file.

    Raises FileExtractionError with a user-facing message when the file is
    empty, corrupt, or a scanned image PDF.
    """
    if file_type == "pdf":
        text = _extract_pdf(path)
    elif file_type == "docx":
        text = _extract_docx(path)
    elif file_type == "doc":
        text = _extract_doc(path)
    elif file_type == "csv":
        text = _extract_csv(path)
    elif file_type == "txt":
        text = _extract_txt(path)
    elif file_type == "xlsx":
        text = _extract_xlsx(path)
    elif file_type == "pptx":
        text = _extract_pptx(path)
    elif file_type in ("jpg", "jpeg", "png"):
        # Health-chatbot policy: never OCR/analyze image content. Tell the
        # user plainly instead of silently dropping the attachment.
        raise FileExtractionError(IMAGE_MESSAGE)
    else:
        raise FileExtractionError(CORRUPT_FILE_MESSAGE)

    if not text or not text.strip():
        raise FileExtractionError(EMPTY_FILE_MESSAGE)
    return _cap_words(text)


if __name__ == "__main__":
    import os
    import tempfile

    print("file_extractor standalone sanity test")
    print("-" * 50)

    with tempfile.TemporaryDirectory() as tmp:
        # txt
        txt_path = os.path.join(tmp, "sample.txt")
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write("Patient has a mild fever and a persistent cough.")
        print("txt type:", resolve_file_type("sample.txt", txt_path))
        print("txt text:", extract_text(txt_path, "txt"))

        # csv
        csv_path = os.path.join(tmp, "log.csv")
        with open(csv_path, "w", encoding="utf-8", newline="") as f:
            w = csv.writer(f)
            w.writerow(["date", "symptom", "severity"])
            w.writerow(["2026-08-01", "fever", "mild"])
            w.writerow(["2026-08-02", "cough", "moderate"])
        print("csv text:\n", extract_text(csv_path, "csv"))

        # mismatched extension -> rejected by magic bytes
        fake_pdf = os.path.join(tmp, "fake.pdf")
        with open(fake_pdf, "wb") as f:
            f.write(b"this is not a real pdf file")
        try:
            resolve_file_type("fake.pdf", fake_pdf)
            print("FAIL: fake pdf accepted")
        except FileExtractionError as exc:
            print("fake pdf rejected:", exc)
